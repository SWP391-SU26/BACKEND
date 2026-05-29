"""
Document loader module - Parse various document formats
Supports: PDF, DOCX, PPTX
"""

from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from typing import List, Tuple
import logging
import os

logger = logging.getLogger(__name__)


class DocumentLoader:
    """
    Load and parse various document formats
    """

    @staticmethod
    def load_pdf(file_path: str) -> str:
        """
        Extract text from PDF file
        
        Args:
            file_path: Path to PDF file
            
        Returns:
            Extracted text
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"PDF file not found: {file_path}")
        
        try:
            text = []
            reader = PdfReader(file_path)
            num_pages = len(reader.pages)
            
            for page_num, page in enumerate(reader.pages):
                page_text = page.extract_text()
                if page_text:
                    text.append(f"--- Page {page_num + 1} ---\n{page_text}")
            
            full_text = "\n".join(text)
            logger.info(f"Loaded PDF: {file_path} ({num_pages} pages)")
            return full_text
        
        except Exception as e:
            logger.error(f"Failed to load PDF {file_path}: {str(e)}")
            raise

    @staticmethod
    def load_docx(file_path: str) -> str:
        """
        Extract text from DOCX file
        
        Args:
            file_path: Path to DOCX file
            
        Returns:
            Extracted text
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"DOCX file not found: {file_path}")
        
        try:
            doc = DocxDocument(file_path)
            text = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    text.append(para.text)
            
            # Extract from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells)
                    if row_text.strip():
                        text.append(row_text)
            
            full_text = "\n".join(text)
            logger.info(f"Loaded DOCX: {file_path} ({len(text)} paragraphs)")
            return full_text
        
        except Exception as e:
            logger.error(f"Failed to load DOCX {file_path}: {str(e)}")
            raise

    @staticmethod
    def load_pptx(file_path: str) -> str:
        """
        Extract text from PPTX file
        
        Args:
            file_path: Path to PPTX file
            
        Returns:
            Extracted text
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"PPTX file not found: {file_path}")
        
        try:
            prs = Presentation(file_path)
            text = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text.append(f"--- Slide {slide_num} ---")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text.append(shape.text)
            
            full_text = "\n".join(text)
            logger.info(f"Loaded PPTX: {file_path} ({len(prs.slides)} slides)")
            return full_text
        
        except Exception as e:
            logger.error(f"Failed to load PPTX {file_path}: {str(e)}")
            raise

    @staticmethod
    def load_file(file_path: str) -> str:
        """
        Auto-detect and load file based on extension
        
        Args:
            file_path: Path to document file
            
        Returns:
            Extracted text
        """
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()
        
        if ext == ".pdf":
            return DocumentLoader.load_pdf(file_path)
        elif ext == ".docx":
            return DocumentLoader.load_docx(file_path)
        elif ext == ".pptx":
            return DocumentLoader.load_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file format: {ext}")

    @staticmethod
    def load_files(file_paths: List[str]) -> List[Tuple[str, str]]:
        """
        Load multiple files
        
        Args:
            file_paths: List of file paths
            
        Returns:
            List of (file_path, text) tuples
        """
        results = []
        
        for file_path in file_paths:
            try:
                text = DocumentLoader.load_file(file_path)
                results.append((file_path, text))
            except Exception as e:
                logger.error(f"Failed to load file {file_path}: {str(e)}")
                # Continue with next file
        
        logger.info(f"Successfully loaded {len(results)}/{len(file_paths)} files")
        return results
 
